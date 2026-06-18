"""Generate JLL Bid Copilot presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io
from PIL import Image, ImageDraw, ImageFont
import os

# ── Brand colours ───────────────────────────────────────────────
RED    = RGBColor(0xE3, 0x06, 0x13)
BLACK  = RGBColor(0x0F, 0x17, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xF1, 0xF5, 0xF9)
DGREY  = RGBColor(0x64, 0x74, 0x8B)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
DARK   = RGBColor(0x14, 0x14, 0x22)
LIGHT_RED = RGBColor(0xFF, 0xF0, 0xF0)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ─────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def add_txt(slide, text, x, y, w, h, size=14, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    return tb

def add_para(tf, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=0, bullet=False):
    from pptx.util import Pt as P
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 1
    r = p.add_run()
    r.text = text
    r.font.size  = P(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    return p

def red_bar(slide, height=Inches(0.55)):
    add_rect(slide, 0, 0, W, height, RED)

def bottom_bar(slide):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), DARK)
    add_txt(slide, "JLL Bid Copilot  |  Internal Use Only  |  Confidential",
            Inches(0.3), H - Inches(0.35), Inches(8), Inches(0.32),
            size=8, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.LEFT)
    add_txt(slide, "© 2026 JLL. All rights reserved.",
            Inches(9), H - Inches(0.35), Inches(4), Inches(0.32),
            size=8, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

def screenshot_box(slide, x, y, w, h, label="[ Screenshot ]", sub=""):
    """Styled placeholder for a screenshot."""
    bg = add_rect(slide, x, y, w, h, RGBColor(0xE8,0xEF,0xF8))
    bg.line.color.rgb = RGBColor(0xCB,0xD5,0xE1)
    bg.line.width = Pt(1)
    # camera icon area
    add_rect(slide, x + w//2 - Inches(0.35), y + h//2 - Inches(0.45),
             Inches(0.7), Inches(0.5), RGBColor(0xCB,0xD5,0xE1))
    add_txt(slide, "📷", x + w//2 - Inches(0.2), y + h//2 - Inches(0.48),
            Inches(0.4), Inches(0.45), size=18, align=PP_ALIGN.CENTER)
    add_txt(slide, label,
            x, y + h//2 + Inches(0.08), w, Inches(0.35),
            size=11, bold=True, color=DGREY, align=PP_ALIGN.CENTER)
    if sub:
        add_txt(slide, sub,
                x, y + h//2 + Inches(0.42), w, Inches(0.35),
                size=9, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)

def stage_pill(slide, x, y, label, color):
    r = add_rect(slide, x, y, Inches(1.45), Inches(0.34), color)
    r.line.fill.background()
    add_txt(slide, label, x, y + Inches(0.04), Inches(1.45), Inches(0.28),
            size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def kpi_box(slide, x, y, w, h, val, lbl, val_color=RED):
    add_rect(slide, x, y, w, h, WHITE).line.color.rgb = RGBColor(0xE2,0xE8,0xF0)
    add_txt(slide, val, x, y + Inches(0.12), w, Inches(0.5),
            size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_txt(slide, lbl, x, y + Inches(0.6), w, Inches(0.3),
            size=9, color=DGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, DARK)
# red accent stripe
add_rect(sl, 0, 0, Inches(0.08), H, RED)
# logo area
logo_box = add_rect(sl, Inches(0.3), Inches(0.28), Inches(1.6), Inches(0.5), RED)
add_txt(sl, "JLL", Inches(0.3), Inches(0.28), Inches(1.6), Inches(0.5),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_txt(sl, "BID COPILOT",
        Inches(0.3), Inches(1.1), Inches(9), Inches(1.2),
        size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(sl, Inches(0.3), Inches(2.4), Inches(3.5), Inches(0.06), RED)

add_txt(sl, "AI-Powered Commercial Real Estate Bidding Assistant",
        Inches(0.3), Inches(2.6), Inches(9), Inches(0.6),
        size=20, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.LEFT)

add_txt(sl, "From requirement intake to FM upsell — the full bid lifecycle, automated.",
        Inches(0.3), Inches(3.3), Inches(9), Inches(0.5),
        size=14, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.LEFT)

# right side decorative boxes
for i, (lbl, col) in enumerate([
    ("Requirements", RGBColor(0x1E,0x3A,0x5F)),
    ("Shortlisting",  RGBColor(0x1A,0x3A,0x4F)),
    ("Proposal",      RGBColor(0x2D,0x1F,0x47)),
    ("Negotiation",   RGBColor(0x1B,0x43,0x32)),
    ("FM Upsell",     RGBColor(0x4A,0x19,0x42)),
]):
    bx = Inches(10.1) + i * Inches(0.62)
    add_rect(sl, bx, Inches(1.5), Inches(0.55), Inches(4.5), col)
    tb = sl.shapes.add_textbox(bx, Inches(3.7), Inches(0.55), Inches(0.5))
    tb.text_frame.text = lbl
    for pr in tb.text_frame.paragraphs:
        pr.alignment = PP_ALIGN.CENTER
        for run in pr.runs:
            run.font.size = Pt(7)
            run.font.color.rgb = RGBColor(0x94,0xA3,0xB8)

add_txt(sl, "Presented by  Neha Sharma, Bid Manager — JLL India",
        Inches(0.3), H - Inches(1.0), Inches(8), Inches(0.35),
        size=11, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.LEFT)
add_txt(sl, "June 2026  |  Confidential",
        Inches(0.3), H - Inches(0.65), Inches(8), Inches(0.35),
        size=10, color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "The Problem with Manual Bidding",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl, "How JLL bid managers work today — and where time is lost",
        Inches(0.4), Inches(0.75), Inches(12), Inches(0.4),
        size=13, color=DGREY)

pain_points = [
    ("⏱  Hours wasted on manual property search",
     "Bid managers sift through spreadsheets to find matching properties. A 250-seat Bengaluru search takes 3–4 hours manually."),
    ("📊  No consistent competitive benchmarking",
     "Each broker has their own intel on CBRE, Colliers, Knight Frank. No shared, up-to-date competitor pricing view."),
    ("📄  Proposal drafting from scratch every time",
     "Every bid proposal is re-written. No template that automatically pulls property data, TCO, and market benchmarks together."),
    ("💡  FM upsell opportunities missed on lost bids",
     "When a lease bid is lost, there is no systematic follow-up for standalone FM services — a significant revenue gap."),
    ("✅  No structured approval workflow",
     "Proposals go to clients without a formal internal review step. No audit trail of who approved what."),
]

for i, (title, body) in enumerate(pain_points):
    y = Inches(1.3) + i * Inches(1.05)
    add_rect(sl, Inches(0.4), y, Inches(0.06), Inches(0.75), RED)
    add_txt(sl, title, Inches(0.65), y, Inches(11.8), Inches(0.38),
            size=13, bold=True, color=BLACK)
    add_txt(sl, body, Inches(0.65), y + Inches(0.38), Inches(11.8), Inches(0.38),
            size=11, color=DGREY)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — What is JLL Bid Copilot
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "What is JLL Bid Copilot?",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl,
    "An internal AI decision-support agent that automates the full commercial real estate bid lifecycle — "
    "from requirement intake to FM upsell — so bid managers close more deals with less manual effort.",
    Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.55),
    size=13, color=BLACK)

# Three column cards
cols = [
    ("🤖  AI-Powered", RED,
     ["Understands natural language", "Extracts requirements from chat",
      "Calls 10 specialised tools", "Works with Claude or DeepSeek"]),
    ("🏢  Real Estate Expert", RGBColor(0x1E,0x40,0xAF),
     ["Lease & purchase transactions", "Bengaluru market data built-in",
      "Competitor intel (CBRE, Colliers…)", "FM pricing by headcount tier"]),
    ("🔒  Human-in-the-Loop", GREEN,
     ["Nothing sent without approval", "Structured approval queue",
      "Full audit trail per opportunity", "Win or lose — find FM revenue"]),
]
for i, (title, col, bullets) in enumerate(cols):
    x = Inches(0.4) + i * Inches(4.3)
    add_rect(sl, x, Inches(1.45), Inches(4.0), Inches(4.8), col)
    add_txt(sl, title, x + Inches(0.15), Inches(1.55), Inches(3.7), Inches(0.5),
            size=14, bold=True, color=WHITE)
    add_rect(sl, x, Inches(2.05), Inches(4.0), Inches(0.03), WHITE)
    for j, b in enumerate(bullets):
        add_txt(sl, f"✓  {b}", x + Inches(0.15),
                Inches(2.2) + j * Inches(0.56), Inches(3.7), Inches(0.48),
                size=12, color=WHITE)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — 10-Stage Workflow
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "The 10-Stage Bid Workflow",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl, "The agent moves through a structured pipeline — every step is tracked, logged, and human-approved before client contact.",
        Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.4), size=12, color=DGREY)

stages = [
    ("1\nINTAKE",           RGBColor(0x1E,0x40,0xAF), "Extracts client\nrequirements"),
    ("2\nSHORTLISTING",     RGBColor(0x05,0x96,0x69), "Searches &\nfilters inventory"),
    ("3\nCOMPARISON",        RGBColor(0x92,0x40,0x0E), "Market & competitor\nbenchmark"),
    ("4\nPROPOSAL",         RED,                       "Drafts bid\nproposal"),
    ("5\nFM BUNDLING",      RGBColor(0x6D,0x28,0xD9), "Configures FM\nservice bundle"),
    ("6\nNEGOTIATION",      RGBColor(0x0E,0x73,0x90), "Talking points\n& strategy"),
    ("7\nOUTCOME",          RGBColor(0x44,0x40,0x3E), "Awaits client\ndecision"),
    ("8\nWON",              GREEN,                     "Closes with FM\ncontract attached"),
    ("9\nFM UPSELL",        RGBColor(0xB4,0x53,0x09), "Standalone FM\npitch on lost bids"),
    ("10\nLOGGED",          RGBColor(0x1F,0x29,0x37), "Archived with\nfull audit trail"),
]

row1, row2 = stages[:5], stages[5:]
for row_idx, row in enumerate([row1, row2]):
    y_base = Inches(1.35) + row_idx * Inches(2.85)
    for i, (label, col, desc) in enumerate(row):
        x = Inches(0.35) + i * Inches(2.58)
        add_rect(sl, x, y_base, Inches(2.3), Inches(1.1), col)
        add_txt(sl, label, x, y_base + Inches(0.1), Inches(2.3), Inches(0.9),
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, x, y_base + Inches(1.1), Inches(2.3), Inches(1.4),
                 RGBColor(0xF8,0xFA,0xFC))
        add_txt(sl, desc, x + Inches(0.1), y_base + Inches(1.18),
                Inches(2.1), Inches(1.2), size=10, color=DGREY)
        # arrow
        if i < 4:
            add_txt(sl, "▶", x + Inches(2.3), y_base + Inches(0.35),
                    Inches(0.28), Inches(0.4), size=11,
                    color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)

# Join row labels
add_txt(sl, "⬇  If bid LOST: agent pivots immediately to FM Upsell (Stage 9) — revenue secured regardless of outcome",
        Inches(0.35), Inches(6.95), Inches(12.5), Inches(0.38),
        size=10, bold=True, color=RGBColor(0xB4,0x53,0x09))
bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Dashboard Overview (screenshot)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Interactive Enterprise Dashboard",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

screenshot_box(sl, Inches(0.35), Inches(0.65), Inches(8.4), Inches(5.8),
               "[ Add dashboard screenshot here ]",
               "http://127.0.0.1:8000  —  full page view with all 13 nav items visible")

# callout labels
callouts = [
    (Inches(9.05), Inches(0.8),  "13-View Navigation",    "Copilot · Dashboard · Opportunities\nRequirements · Property Search\nComparisons · Proposals · FM Services\nWin/Loss · Pipeline · Reports…"),
    (Inches(9.05), Inches(2.55), "AI Chat Panel",         "Natural language interface.\nType requirements, get\nshortlists and proposals."),
    (Inches(9.05), Inches(4.05), "FM Services Bundle",    "Live checkbox configurator.\n₹4,45,000/month default.\nOne-click add to proposal."),
    (Inches(9.05), Inches(5.35), "Approval Queue",        "Inline Approve / Reject cards.\nNothing sent to client\nwithout human sign-off."),
]
for x, y, title, body in callouts:
    add_rect(sl, x, y, Inches(0.05), Inches(0.9), RED)
    add_txt(sl, title, x + Inches(0.15), y, Inches(3.85), Inches(0.3),
            size=11, bold=True, color=BLACK)
    add_txt(sl, body, x + Inches(0.15), y + Inches(0.3), Inches(3.85), Inches(0.65),
            size=9.5, color=DGREY)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Run AI Demo
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Live AI Workflow Demo",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl, 'Click "Run AI Demo" in the dashboard header to watch the agent process a real bid in real time.',
        Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.4), size=12, color=DGREY)

# Left: screenshot placeholder
screenshot_box(sl, Inches(0.35), Inches(1.2), Inches(7.0), Inches(5.4),
               "[ Add demo modal screenshot here ]",
               "Dark two-column modal — agent log on left, output cards on right")

# Right: 10 steps list
add_txt(sl, "What the demo shows:", Inches(7.7), Inches(1.2), Inches(5.3), Inches(0.4),
        size=13, bold=True, color=BLACK)

steps = [
    ("1", "Bid created",            "OPP-2026-0089 initialised"),
    ("2", "Requirements extracted", "250 seats · ₹130–150 · Q3 2026"),
    ("3", "Property search",        "23 candidates → 5 shortlisted"),
    ("4", "Shortlist scored",       "Prestige Tech Park: 92% fit"),
    ("5", "Market benchmark",       "Avg ₹142 · JLL bid 2.8% below"),
    ("6", "Competitor intel",       "CBRE ₹145 HIGH · KF ₹148 LOW"),
    ("7", "Proposal drafted",       "₹138/sqft · ₹8.28Cr/yr · 73% win"),
    ("8", "Approval requested",     "⏸  Workflow paused for Neha"),
    ("9", "FM bundle ready",        "₹4,45,000/month · ₹53.4L/yr"),
    ("10","Demo complete",          "Bid ready for approval"),
]
for i, (num, title, sub) in enumerate(steps):
    y = Inches(1.7) + i * Inches(0.5)
    add_rect(sl, Inches(7.7), y + Inches(0.04), Inches(0.3), Inches(0.3),
             RED if i < 9 else GREEN)
    add_txt(sl, num, Inches(7.7), y + Inches(0.04), Inches(0.3), Inches(0.3),
            size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, title, Inches(8.12), y, Inches(2.5), Inches(0.3),
            size=10, bold=True, color=BLACK)
    add_txt(sl, sub, Inches(10.65), y, Inches(2.5), Inches(0.3),
            size=9, color=DGREY)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Property Shortlisting & Scoring
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "AI Property Shortlisting & Fit Scoring",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl, "ABC Tech Pvt. Ltd. — 250 seats, Bengaluru ORR/Whitefield, ₹130–150/sqft, Q3 2026",
        Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.38), size=12, color=DGREY)

# Table header
hdr_cols = ["Property", "Location", "Rent /sqft", "Seats", "Available", "Fit Score", "AI Verdict"]
col_widths = [2.3, 2.1, 1.1, 0.8, 1.0, 0.9, 1.8]
x0 = Inches(0.35)
y0 = Inches(1.2)
add_rect(sl, x0, y0, Inches(12.6), Inches(0.42), DARK)
cx = x0
for col, cw in zip(hdr_cols, col_widths):
    add_txt(sl, col, cx + Inches(0.08), y0 + Inches(0.06),
            Inches(cw), Inches(0.3), size=9.5, bold=True, color=WHITE)
    cx += Inches(cw)

props = [
    ("Prestige Tech Park",       "Whitefield",    "₹138", "280", "Jul 2026", "92%", "✅ TOP PICK",      GREEN,                     True),
    ("Embassy Manyata BP",       "Outer Ring Rd", "₹143", "340", "Aug 2026", "88%", "✅ Recommended",   GREEN,                     False),
    ("RMZ Ecoworld",             "Bellandur",     "₹135", "255", "Jun 2026", "85%", "✅ Best value",    GREEN,                     False),
    ("Bagmane Tech Park",        "CV Raman Nagar","₹128", "310", "Sep 2026", "78%", "⚠  Pantry gap",   AMBER,                     False),
    ("Brigade Tech Garden",      "Whitefield",    "₹141", "270", "Oct 2026", "74%", "⚠  Timeline risk",AMBER,                     False),
]
for i, (name, loc, rent, seats, avail, fit, verdict, vcol, highlight) in enumerate(props):
    y = y0 + Inches(0.42) + i * Inches(0.92)
    bg = LIGHT_RED if highlight else (GREY if i % 2 == 0 else WHITE)
    add_rect(sl, x0, y, Inches(12.6), Inches(0.88), bg)
    if highlight:
        add_rect(sl, x0, y, Inches(0.06), Inches(0.88), RED)
    row_data = [name, loc, rent, seats, avail, fit, verdict]
    cx = x0
    for j, (val, cw) in enumerate(zip(row_data, col_widths)):
        fc = vcol if j == 6 else (RED if j == 2 else BLACK)
        fb = highlight and j == 0
        add_txt(sl, val, cx + Inches(0.1), y + Inches(0.27),
                Inches(cw), Inches(0.35), size=10.5, bold=fb, color=fc)
        cx += Inches(cw)

add_txt(sl, "AI scores 8 dimensions: Budget fit · Location · Facilities · Timeline · Seats · Grade · Expansion potential · Landlord risk",
        x0, Inches(6.95), Inches(12.6), Inches(0.32), size=9, color=DGREY)
bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Competitive Intelligence
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Competitive Intelligence & Market Positioning",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

# Left: competitor table
add_txt(sl, "Competitor Overview — Bengaluru ORR/Whitefield",
        Inches(0.4), Inches(0.72), Inches(7.5), Inches(0.38), size=12, bold=True, color=BLACK)

comp_hdrs = ["Broker", "Quoted Rent", "Risk Level", "JLL Advantage"]
comp_cw   = [1.7, 1.3, 1.1, 3.5]
y0 = Inches(1.18)
add_rect(sl, Inches(0.35), y0, Inches(7.85), Inches(0.38), DARK)
cx = Inches(0.35)
for h, cw in zip(comp_hdrs, comp_cw):
    add_txt(sl, h, cx+Inches(0.08), y0+Inches(0.05), Inches(cw), Inches(0.28),
            size=9.5, bold=True, color=WHITE)
    cx += Inches(cw)

comps = [
    ("CBRE",              "₹145/sqft", "HIGH",   RED,   "No FM · Incumbent in ORR · JLL 4.9% cheaper"),
    ("Cushman & Wakefield","₹146/sqft", "MEDIUM", AMBER, "FM separate contract · 10–15% pricier FM"),
    ("Colliers",          "₹140/sqft", "MEDIUM", AMBER, "No FM · Flex lease only · Smaller network"),
    ("Knight Frank",      "₹148/sqft", "LOW",    GREEN, "Premium segment · Outside budget ceiling"),
    ("JLL (Our Bid)",     "₹138/sqft", "OUR BID",GREEN, "Best value + integrated FM bundle"),
]
for i, (name, rent, risk, rcol, adv) in enumerate(comps):
    y = y0 + Inches(0.38) + i * Inches(0.82)
    bg = LIGHT_RED if name.startswith("JLL") else (GREY if i%2==0 else WHITE)
    add_rect(sl, Inches(0.35), y, Inches(7.85), Inches(0.78), bg)
    if name.startswith("JLL"):
        add_rect(sl, Inches(0.35), y, Inches(0.06), Inches(0.78), RED)
    vals = [name, rent, risk, adv]
    cx = Inches(0.35)
    for j, (v, cw) in enumerate(zip(vals, comp_cw)):
        fc = rcol if j == 2 else (RED if name.startswith("JLL") and j==0 else BLACK)
        fb = name.startswith("JLL")
        add_txt(sl, v, cx+Inches(0.1), y+Inches(0.22), Inches(cw), Inches(0.35),
                size=10, bold=fb, color=fc)
        cx += Inches(cw)

# Right: market stats
add_txt(sl, "Market Benchmark",
        Inches(8.5), Inches(0.72), Inches(4.5), Inches(0.38), size=12, bold=True, color=BLACK)

kpis_right = [
    ("₹142", "Avg. Market Rent\n/sqft (All Incl.)", DGREY),
    ("₹138", "JLL Bid — Prestige\n(2.8% below mkt)", RED),
    ("17.2%","Vacancy Rate\nBengaluru ORR",         DGREY),
    ("+4.8%","YoY Rent Growth",                     GREEN),
]
for i, (val, lbl, vc) in enumerate(kpis_right):
    kx = Inches(8.5) + (i % 2) * Inches(2.3)
    ky = Inches(1.3) + (i // 2) * Inches(1.5)
    kpi_box(sl, kx, ky, Inches(2.0), Inches(1.25), val, lbl, vc)

add_txt(sl, "Key JLL differentiator",
        Inches(8.5), Inches(4.5), Inches(4.5), Inches(0.35), size=11, bold=True, color=BLACK)
add_rect(sl, Inches(8.5), Inches(4.92), Inches(4.5), Inches(1.5),
         RGBColor(0xF0,0xFD,0xF4))
add_txt(sl,
    "Integrated leasing + FM bundle.\n\n"
    "Saves ABC Tech ₹6–8L/year vs. sourcing FM separately.\n"
    "Single JLL contract vs. two vendors.",
    Inches(8.65), Inches(5.0), Inches(4.2), Inches(1.3),
    size=11, color=RGBColor(0x06,0x5F,0x46))

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — New Bid Wizard
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Creating a New Bid — 4-Step Wizard",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

add_txt(sl, 'Click "+ New Bid" in the top-right of the dashboard to open the guided wizard.',
        Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.38), size=12, color=DGREY)

# Step boxes
step_data = [
    ("Step 1", "Client Details", RED,
     ["Company name", "Industry sector", "Primary contact", "Deal type"]),
    ("Step 2", "Location & Space", RGBColor(0x1E,0x40,0xAF),
     ["City & sub-market", "Preferred locations", "Seats required", "Sq.ft range"]),
    ("Step 3", "Commercial Terms", RGBColor(0x05,0x96,0x69),
     ["Budget range (₹/sqft)", "Lease term (months)", "Move-in timeline", "Competitor pricing"]),
    ("Step 4", "FM Services", RGBColor(0x6D,0x28,0xD9),
     ["Leasing only / Leasing+FM", "Select FM services", "Live cost estimate", "Submit & Analyse"]),
]
for i, (step, title, col, items) in enumerate(step_data):
    x = Inches(0.35) + i * Inches(3.22)
    add_rect(sl, x, Inches(1.25), Inches(3.0), Inches(0.55), col)
    add_txt(sl, step, x + Inches(0.12), Inches(1.3), Inches(1.0), Inches(0.42),
            size=10, bold=True, color=WHITE)
    add_txt(sl, title, x + Inches(1.1), Inches(1.3), Inches(1.8), Inches(0.42),
            size=13, bold=True, color=WHITE)
    add_rect(sl, x, Inches(1.8), Inches(3.0), Inches(3.8), GREY)
    for j, item in enumerate(items):
        add_txt(sl, f"▸  {item}",
                x + Inches(0.18), Inches(2.05) + j * Inches(0.72),
                Inches(2.7), Inches(0.55), size=11.5, color=BLACK)
    # connector arrow
    if i < 3:
        add_txt(sl, "→", x + Inches(3.0), Inches(2.55),
                Inches(0.22), Inches(0.4), size=18,
                color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)

# Bottom result box
add_rect(sl, Inches(0.35), Inches(5.8), Inches(12.6), Inches(0.75),
         RGBColor(0xF0,0xFD,0xF4))
add_rect(sl, Inches(0.35), Inches(5.8), Inches(0.06), Inches(0.75), GREEN)
add_txt(sl,
    "After submission → AI analyses feasibility, searches inventory, scores shortlist, "
    "and auto-populates the full dashboard with the new opportunity in under 30 seconds.",
    Inches(0.6), Inches(5.92), Inches(12.1), Inches(0.52),
    size=12, color=RGBColor(0x06,0x5F,0x46))

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — FM Services Upsell
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "FM Services — Win or Lose, Find Revenue",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

# Left side
add_txt(sl, "Recommended FM Bundle — 250 Seats",
        Inches(0.4), Inches(0.75), Inches(6.5), Inches(0.38), size=13, bold=True, color=BLACK)

fm_services = [
    ("☕", "Pantry Management",     "₹1,20,000", RGBColor(0xFF,0xF3,0xE0)),
    ("⚡", "Electrical Maintenance","₹75,000",   RGBColor(0xFF,0xFD,0xE7)),
    ("⬥", "Washroom Services",     "₹60,000",   RGBColor(0xE3,0xF2,0xFD)),
    ("⚙", "Plumbing Services",     "₹50,000",   RGBColor(0xF3,0xE5,0xF5)),
    ("👥", "General Ops Support",  "₹1,00,000", RGBColor(0xFC,0xE4,0xEC)),
    ("✎", "Stationery & Supplies", "₹40,000",   RGBColor(0xE8,0xF5,0xE9)),
]
for i, (ico, name, amt, bg) in enumerate(fm_services):
    y = Inches(1.25) + i * Inches(0.78)
    add_rect(sl, Inches(0.4), y, Inches(6.0), Inches(0.66), bg)
    add_txt(sl, ico,  Inches(0.55), y + Inches(0.16), Inches(0.38), Inches(0.38), size=16)
    add_txt(sl, name, Inches(1.05), y + Inches(0.18), Inches(3.8), Inches(0.35),
            size=12, bold=True, color=BLACK)
    add_txt(sl, amt,  Inches(4.9), y + Inches(0.18), Inches(1.4), Inches(0.35),
            size=12, color=RED, align=PP_ALIGN.RIGHT)

add_rect(sl, Inches(0.4), Inches(6.0), Inches(6.0), Inches(0.6),
         RGBColor(0x0F,0x17,0x2A))
add_txt(sl, "Monthly FM Total:",
        Inches(0.6), Inches(6.08), Inches(3.0), Inches(0.42),
        size=13, bold=True, color=WHITE)
add_txt(sl, "₹4,45,000 / month",
        Inches(3.5), Inches(6.05), Inches(2.8), Inches(0.5),
        size=18, bold=True, color=RED, align=PP_ALIGN.RIGHT)

# Right side
add_txt(sl, "Why FM Upsell Matters",
        Inches(7.0), Inches(0.75), Inches(6.0), Inches(0.38), size=13, bold=True, color=BLACK)

fm_points = [
    ("₹53.4L/year", "Annual FM revenue per client"),
    ("Win or lose", "FM pitch fires automatically on lost bids"),
    ("Single contract", "JLL manages leasing + FM — one relationship"),
    ("6–8L savings", "Vs. client sourcing FM separately"),
]
for i, (val, lbl) in enumerate(fm_points):
    y = Inches(1.25) + i * Inches(1.0)
    add_rect(sl, Inches(7.0), y, Inches(5.9), Inches(0.85),
             LIGHT_RED if i == 0 else GREY)
    add_txt(sl, val, Inches(7.15), y + Inches(0.08), Inches(2.5), Inches(0.42),
            size=20, bold=True, color=RED)
    add_txt(sl, lbl, Inches(7.15), y + Inches(0.48), Inches(5.5), Inches(0.3),
            size=11, color=DGREY)

add_rect(sl, Inches(7.0), Inches(5.4), Inches(5.9), Inches(1.22),
         RGBColor(0x1A,0x3A,0x4F))
add_txt(sl,
    '"Wherever you land, JLL runs the facility."\n\n'
    'The agent auto-prepares the standalone FM pitch the moment a lease bid is lost.',
    Inches(7.2), Inches(5.5), Inches(5.5), Inches(1.0),
    size=11, color=WHITE)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Benefits / ROI
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Business Impact & Key Benefits",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

metrics = [
    ("3–4 hrs → 30 sec",  "Property shortlisting time",   RED),
    ("73%",               "Win probability — Prestige bid",GREEN),
    ("₹53.4L/yr",         "FM revenue per 250-seat client",RED),
    ("0",                 "Proposals sent without approval",GREEN),
    ("10 stages",         "Fully automated workflow steps", RED),
    ("2 LLMs",            "Anthropic or DeepSeek — your choice", DGREY),
]
for i, (val, lbl, col) in enumerate(metrics):
    x = Inches(0.4) + (i % 3) * Inches(4.25)
    y = Inches(0.95) + (i // 3) * Inches(1.85)
    add_rect(sl, x, y, Inches(3.9), Inches(1.6), GREY)
    add_rect(sl, x, y, Inches(3.9), Inches(0.08), col)
    add_txt(sl, val, x, y + Inches(0.18), Inches(3.9), Inches(0.7),
            size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_txt(sl, lbl, x, y + Inches(0.9), Inches(3.9), Inches(0.5),
            size=11, color=DGREY, align=PP_ALIGN.CENTER)

# Bottom summary
add_rect(sl, Inches(0.4), Inches(5.55), Inches(12.5), Inches(1.05),
         RGBColor(0x0F,0x17,0x2A))
benefits = [
    "✓  No manual spreadsheet work",
    "✓  Competitor intel always current",
    "✓  Proposals consistent & data-backed",
    "✓  FM revenue captured on every deal",
    "✓  Full audit trail for compliance",
]
for i, b in enumerate(benefits):
    add_txt(sl, b, Inches(0.7) + i * Inches(2.44), Inches(5.75),
            Inches(2.4), Inches(0.65), size=10.5, bold=True, color=WHITE)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — How to Use (Step by Step)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "How to Use — Step by Step",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

how_steps = [
    ("1", RED,                     "Open the Dashboard",
     'Navigate to http://127.0.0.1:8000\nThe ABC Tech demo loads automatically. '
     'Explore all 13 navigation views on the left sidebar.'),
    ("2", RGBColor(0x1E,0x40,0xAF),"Run AI Demo",
     'Click "Run AI Demo" (top-right). Watch the animated 10-step workflow — '
     'tool calls on the left, output cards on the right.'),
    ("3", RGBColor(0x05,0x96,0x69),"Create a New Bid",
     'Click "+ New Bid". Fill in Client Details → Location & Space → '
     'Commercial Terms → FM Services. Submit — AI populates the dashboard instantly.'),
    ("4", RGBColor(0x6D,0x28,0xD9),"Chat with the Copilot",
     'Use the right-side chat panel. Ask: "Find the best properties", '
     '"Compare competitors", "Draft a proposal", "What is our negotiation strategy?"'),
    ("5", RGBColor(0x92,0x40,0x0E),"Review & Approve",
     'When the agent raises an approval, an Approve / Reject card appears in chat. '
     'Nothing goes to the client until you click Approve.'),
    ("6", GREEN,                   "Close the Bid",
     'Mark as Won → FM contract auto-attached. Mark as Lost → FM standalone pitch '
     'prepared automatically. Win or lose — JLL captures revenue.'),
]
for i, (num, col, title, body) in enumerate(how_steps):
    row, c = divmod(i, 2)
    x = Inches(0.35) + c * Inches(6.45)
    y = Inches(1.05) + row * Inches(1.95)
    add_rect(sl, x, y, Inches(0.52), Inches(0.52), col)
    add_txt(sl, num, x, y, Inches(0.52), Inches(0.52),
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, title, x + Inches(0.65), y, Inches(5.5), Inches(0.42),
            size=13, bold=True, color=BLACK)
    add_txt(sl, body, x + Inches(0.65), y + Inches(0.44), Inches(5.5), Inches(1.2),
            size=10.5, color=DGREY)

bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Technical Setup
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
red_bar(sl)
add_txt(sl, "JLL", Inches(0.28), Inches(0.08), Inches(0.8), Inches(0.38),
        size=16, bold=True, color=WHITE)
add_txt(sl, "Technical Setup — Get Running in 5 Minutes",
        Inches(1.2), Inches(0.08), Inches(9), Inches(0.38),
        size=18, bold=True, color=WHITE)

# Left: setup steps
setup = [
    ("Prerequisites", "Python 3.10+  ·  Port 8000 available"),
    ("1  Clone / download",
     "Download ZIP from github.com/definedata234-ux/jll-bidding\nor: git clone https://github.com/definedata234-ux/jll-bidding.git"),
    ("2  Create virtual env",
     "python -m venv .venv\n.venv\\Scripts\\activate"),
    ("3  Install dependencies",
     "pip install -r requirements.txt"),
    ("4  Configure API key",
     "copy .env.example .env\n# Edit .env → set DEEPSEEK_API_KEY or ANTHROPIC_API_KEY"),
    ("5  Start server",
     "uvicorn app.api:app --reload --port 8000"),
    ("6  Open browser",
     "http://127.0.0.1:8000   ← full demo loads automatically"),
]
for i, (title, code) in enumerate(setup):
    y = Inches(0.85) + i * Inches(0.85)
    is_cmd = i > 0
    add_rect(sl, Inches(0.35), y, Inches(0.06),
             Inches(0.7 if '\n' in code else 0.5), RED if is_cmd else DGREY)
    add_txt(sl, title, Inches(0.55), y, Inches(7.0), Inches(0.32),
            size=11, bold=True, color=BLACK)
    add_rect(sl, Inches(0.55), y + Inches(0.32), Inches(7.0),
             Inches(0.45 if '\n' in code else 0.32), RGBColor(0xF1,0xF5,0xF9))
    add_txt(sl, code, Inches(0.7), y + Inches(0.33), Inches(6.8),
            Inches(0.42 if '\n' in code else 0.28),
            size=9.5, color=RGBColor(0x1E,0x40,0xAF))

# Right: architecture
add_txt(sl, "Architecture",
        Inches(8.0), Inches(0.85), Inches(5.0), Inches(0.38), size=13, bold=True, color=BLACK)

arch = [
    ("Frontend",   "Vanilla JS · index.html · styles.css · app.js",    RGBColor(0x1E,0x40,0xAF)),
    ("API Layer",  "FastAPI · uvicorn · slowapi rate limiting",          RGBColor(0x05,0x96,0x69)),
    ("Agent",      "Orchestrator → 10 tools → Pydantic validation",     RED),
    ("Services",   "Property search · Scoring · FM pricing",             RGBColor(0x6D,0x28,0xD9)),
    ("Data",       "Repositories → mock JSON (swap for live DB)",        DGREY),
    ("LLM",        "DeepSeek (default) or Anthropic Claude",            RGBColor(0x92,0x40,0x0E)),
]
for i, (layer, detail, col) in enumerate(arch):
    y = Inches(1.35) + i * Inches(0.82)
    add_rect(sl, Inches(8.0), y, Inches(1.3), Inches(0.62), col)
    add_txt(sl, layer, Inches(8.0), y + Inches(0.12), Inches(1.3), Inches(0.38),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(9.35), y, Inches(3.8), Inches(0.62), GREY)
    add_txt(sl, detail, Inches(9.5), y + Inches(0.14), Inches(3.6), Inches(0.35),
            size=9.5, color=BLACK)
    if i < 5:
        add_txt(sl, "↓", Inches(8.55), y + Inches(0.62), Inches(0.3), Inches(0.2),
                size=9, color=DGREY, align=PP_ALIGN.CENTER)

add_txt(sl, "Demo mode works without any API key — full UI and all 13 views load from built-in demo data.",
        Inches(8.0), Inches(6.52), Inches(5.1), Inches(0.4),
        size=9.5, bold=True, color=GREEN)
bottom_bar(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Closing / Next Steps
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, DARK)
add_rect(sl, 0, 0, Inches(0.08), H, RED)
add_txt(sl, "JLL", Inches(0.3), Inches(0.28), Inches(1.6), Inches(0.5),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_txt(sl, "From First Contact\nto FM Contract.",
        Inches(0.4), Inches(1.2), Inches(8), Inches(2.0),
        size=44, bold=True, color=WHITE)

add_rect(sl, Inches(0.4), Inches(3.3), Inches(3.5), Inches(0.06), RED)

add_txt(sl, "JLL Bid Copilot automates the bid lifecycle so your team\nfocuses on relationships — not spreadsheets.",
        Inches(0.4), Inches(3.5), Inches(8.5), Inches(0.9),
        size=15, color=RGBColor(0xCB,0xD5,0xE1))

next_steps = [
    "🔗  Live demo: http://127.0.0.1:8000",
    "📂  Source code: github.com/definedata234-ux/jll-bidding",
    "📧  Questions: Contact Neha Sharma, Bid Manager",
]
for i, s in enumerate(next_steps):
    add_txt(sl, s, Inches(0.4), Inches(4.65) + i * Inches(0.55),
            Inches(8.5), Inches(0.45), size=13, color=RGBColor(0x94,0xA3,0xB8))

# Right large stat
add_txt(sl, "₹4,45,000",
        Inches(9.5), Inches(1.8), Inches(3.5), Inches(1.0),
        size=40, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_txt(sl, "monthly FM revenue\nper client — win or lose",
        Inches(9.5), Inches(2.8), Inches(3.5), Inches(0.8),
        size=13, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)

add_txt(sl, "73%",
        Inches(9.5), Inches(3.9), Inches(3.5), Inches(0.85),
        size=40, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_txt(sl, "predicted win probability\non top-scored opportunity",
        Inches(9.5), Inches(4.75), Inches(3.5), Inches(0.7),
        size=13, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)

add_txt(sl, "© 2026 JLL. All rights reserved.  |  Confidential",
        Inches(0.3), H - Inches(0.45), Inches(8), Inches(0.35),
        size=9, color=RGBColor(0x47,0x55,0x69))


# ── Save ─────────────────────────────────────────────────────────
out = r"c:\Users\sovan\OneDrive\Desktop\JLL_Bid_Copilot_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
